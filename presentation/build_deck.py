"""Generate the MoveInSync Pulse hackathon deck.

    python presentation/build_deck.py

Produces presentation/MoveInSync_Pulse_Hackathon.pptx (9 slides, 16:9).
Enterprise SaaS visual language: white background, navy text, MoveInSync blue,
teal accent, sparing status colors, minimal text. Native shapes only — no
external images required. Drop real screenshots from docs/screenshots/ onto
slide 7 after generation.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

NAVY = RGBColor(0x0F, 0x25, 0x40)
BLUE = RGBColor(0x16, 0x68, 0xE3)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
BG = RGBColor(0xF7, 0xF9, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x5B, 0x6B, 0x80)
LIGHT = RGBColor(0xEE, 0xF4, 0xFB)
BORDER = RGBColor(0xD5, 0xDF, 0xEC)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x0E, 0x9F, 0x6E)

EMU = 914400
SW, SH = int(13.333 * EMU), int(7.5 * EMU)


def inches(v: float) -> int:
    return int(v * EMU)


def bg(slide) -> None:
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    slide.shapes._spTree.remove(r._element)
    slide.shapes._spTree.insert(2, r._element)


def box(slide, x, y, w, h, *, fill=WHITE, line=BORDER, radius=True, line_w=1.0):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        inches(x), inches(y), inches(w), inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    if isinstance(runs, tuple):
        runs = [runs]
    for i, (s, size, color, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = s
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = "Segoe UI"
    return tb


def chip(slide, x, y, w, label, *, fill=LIGHT, fg=NAVY, line=None, size=12):
    b = box(slide, x, y, w, 0.42, fill=fill, line=line)
    tf = b.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(size); r.font.color.rgb = fg; r.font.bold = True
    r.font.name = "Segoe UI"
    return b


def arrow_down(slide, x, y):
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, inches(x), inches(y), inches(0.3), inches(0.28))
    a.fill.solid(); a.fill.fore_color.rgb = BLUE; a.line.fill.background()
    a.shadow.inherit = False


def eyebrow(slide, x, y, label, color=BLUE):
    text(slide, x, y, 8, 0.3, (label.upper(), 12, color, True))


def title(slide, x, y, label, size=30):
    text(slide, x, y, 11.5, 0.9, (label, size, NAVY, True))


prs = Presentation()
prs.slide_width, prs.slide_height = Emu(SW), Emu(SH)
blank = prs.slide_layouts[6]


def new():
    s = prs.slides.add_slide(blank); bg(s); return s


# 1 — Title
s = new()
box(s, 0, 0, 13.333, 7.5, fill=NAVY, line=None, radius=False)
chip(s, 0.9, 1.4, 2.2, "MOVEINSYNC", fill=BLUE, fg=WHITE, size=12)
text(s, 0.9, 2.3, 11.5, 1.6, ("MoveInSync Pulse", 54, WHITE, True))
text(s, 0.9, 3.7, 11.5, 0.8, ("Autonomous Cross-Signal Mobility Intelligence", 24, RGBColor(0xB9,0xCD,0xF0), False))
box(s, 0.9, 4.9, 5.6, 0.9, fill=RGBColor(0x14,0x2F,0x52), line=None)
text(s, 1.15, 5.0, 5.2, 0.7, ('"Detect what dashboards miss."', 20, TEAL, True), anchor=MSO_ANCHOR.MIDDLE)

# 2 — Problem
s = new()
eyebrow(s, 0.9, 0.6, "The problem")
title(s, 0.9, 0.95, "Mobility data is everywhere. Insight isn't.")
domains = ["Trips", "Employees", "Safety", "Billing", "Feedback"]
for i, d in enumerate(domains):
    chip(s, 0.9 + i * 2.35, 2.1, 2.1, d, fill=LIGHT, fg=NAVY, line=BORDER, size=15)
flow = ["Data-rich", "Report-heavy", "Manual correlation", "Slow action"]
for i, f in enumerate(flow):
    b = box(s, 0.9 + i * 3.0, 3.2, 2.5, 0.7, fill=WHITE, line=BORDER)
    b.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = b.text_frame.paragraphs[0].add_run(); r.text = f
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = NAVY; r.font.name = "Segoe UI"
    if i < 3:
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, inches(3.42 + i * 3.0), inches(3.37), inches(0.42), inches(0.36))
        a.fill.solid(); a.fill.fore_color.rgb = SLATE; a.line.fill.background(); a.shadow.inherit = False
box(s, 0.9, 4.7, 11.5, 1.4, fill=NAVY, line=None)
text(s, 1.2, 4.95, 11.0, 1.0,
     ('"The anomaly may not exist in any one report. It emerges when signals are combined."',
      22, WHITE, True), anchor=MSO_ANCHOR.MIDDLE)

# 3 — Solution
s = new()
eyebrow(s, 0.9, 0.6, "Our solution")
title(s, 0.9, 0.95, "Sense → Correlate → Investigate → Reason → Act")
steps = [("Sense", "deterministic\nanalytics"), ("Correlate", "cross-domain\nanomaly engine"),
         ("Investigate", "agent gathers\nevidence"), ("Reason", "LLM explains\nevidence"),
         ("Act", "grounded\nrecommendation")]
for i, (h, sub) in enumerate(steps):
    x = 0.7 + i * 2.5
    b = box(s, x, 2.4, 2.2, 1.6, fill=WHITE, line=BLUE, line_w=1.5)
    text(s, x + 0.1, 2.6, 2.0, 0.5, (h, 17, BLUE, True), align=PP_ALIGN.CENTER)
    text(s, x + 0.1, 3.15, 2.0, 0.8, (sub, 12, SLATE, False), align=PP_ALIGN.CENTER)
    if i < 4:
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, inches(x + 2.22), inches(3.0), inches(0.32), inches(0.3))
        a.fill.solid(); a.fill.fore_color.rgb = TEAL; a.line.fill.background(); a.shadow.inherit = False
box(s, 0.9, 4.9, 11.5, 1.2, fill=LIGHT, line=BORDER)
text(s, 1.2, 5.05, 11.0, 0.9,
     ("Pulse continuously correlates mobility operations data, investigates the anomalies it finds, "
      "explains their business significance, and recommends the next action — every number grounded in analytics.",
      16, NAVY, False), anchor=MSO_ANCHOR.MIDDLE)

# 4 — Cross-signal intelligence (real anomaly)
s = new()
eyebrow(s, 0.9, 0.6, "Cross-signal intelligence")
title(s, 0.9, 0.95, "One vendor. Signals pointing opposite ways.")
card = box(s, 0.9, 2.1, 6.0, 3.9, fill=WHITE, line=BORDER)
chip(s, 1.2, 2.35, 1.2, "HIGH", fill=RGBColor(0xFE,0xE2,0xE2), fg=RED, size=12)
text(s, 5.2, 2.35, 1.5, 0.4, ("risk 91", 15, NAVY, True), align=PP_ALIGN.RIGHT)
text(s, 1.2, 2.95, 5.4, 0.5, ("Aarav Petrov Travel", 20, NAVY, True))
text(s, 1.2, 3.4, 5.4, 0.4, ("Safety divergence", 15, SLATE, False))
chip(s, 1.2, 4.0, 2.6, "Safety alerts  ↑", fill=RGBColor(0xFE,0xE2,0xE2), fg=RED, size=13)
chip(s, 3.95, 4.0, 2.6, "No-show  ↓ better", fill=RGBColor(0xD1,0xFA,0xE5), fg=GREEN, size=13)
chip(s, 1.2, 4.55, 2.6, "Delay  ↓ better", fill=RGBColor(0xD1,0xFA,0xE5), fg=GREEN, size=13)
chip(s, 3.95, 4.55, 2.6, "vs peer median ↑", fill=RGBColor(0xFE,0xE2,0xE2), fg=RED, size=13)
text(s, 1.2, 5.2, 5.4, 0.7,
     ("Alerts 139/1k (~+55% vs June; peer 68.6) while service metrics improved.",
      13, SLATE, False))
text(s, 7.3, 2.3, 5.1, 0.5, ("Why it matters", 18, BLUE, True))
text(s, 7.3, 2.9, 5.1, 3.0,
     [("A single vendor-health score averages these movements together — and misses that the "
       "problem is safety-specific, not general decline.", 16, NAVY, False),
      ("", 8, NAVY, False),
      ("This pattern only appears when domains are correlated. Billing risks are flagged the same "
       "conservative way: potential irregularities requiring reconciliation review, never proven fraud.",
       15, SLATE, False)])

# 5 — How it works
s = new()
eyebrow(s, 0.9, 0.6, "How it works")
title(s, 0.9, 0.95, "Deterministic evidence, grounded narrative")
pipe = ["CSV sources", "Data quality &\nnormalization", "DuckDB +\nParquet",
        "Deterministic\nanalytics", "Cross-domain\nanomaly engine",
        "LangGraph\norchestrator", "LLM synthesis\n(grounded)", "FastAPI +\nReact UI"]
for i, p in enumerate(pipe):
    row, col = divmod(i, 4)
    x = 0.9 + col * 3.0
    y = 2.2 + row * 1.5
    hl = i in (3, 4)
    b = box(s, x, y, 2.6, 1.1, fill=(LIGHT if hl else WHITE), line=(BLUE if hl else BORDER),
            line_w=1.5 if hl else 1.0)
    text(s, x + 0.1, y + 0.15, 2.4, 0.8, (p, 13, (BLUE if hl else NAVY), hl), align=PP_ALIGN.CENTER)
    if col < 3 and i < 7:
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, inches(x + 2.62), inches(y + 0.42), inches(0.3), inches(0.28))
        a.fill.solid(); a.fill.fore_color.rgb = SLATE; a.line.fill.background(); a.shadow.inherit = False
box(s, 0.9, 5.55, 11.5, 0.95, fill=RGBColor(0xFF,0xF7,0xED), line=AMBER, line_w=1.2)
text(s, 1.2, 5.62, 11.0, 0.8,
     ("⚠  The LLM never calculates operational metrics — it only explains supplied evidence. "
      "Raw trip data never reaches the model.", 16, RGBColor(0x92,0x40,0x0E), True),
     anchor=MSO_ANCHOR.MIDDLE)

# 6 — Agentic investigation
s = new()
eyebrow(s, 0.9, 0.6, "Agentic investigation")
title(s, 0.9, 0.95, "The agent decides what to investigate")
chain = ["Detected anomaly", "Agent chooses tools", "Vendor analytics",
         "Safety analytics", "Delay context", "Data quality", "Grounded recommendation"]
for i, c in enumerate(chain):
    y = 1.9 + i * 0.68
    accent = i in (0, 6)
    b = box(s, 3.6, y, 6.1, 0.56, fill=(NAVY if accent else WHITE), line=(None if accent else BORDER))
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = c
    r.font.size = Pt(14); r.font.bold = True; r.font.name = "Segoe UI"
    r.font.color.rgb = WHITE if accent else NAVY
    if i < 6:
        arrow_down(s, 6.5, y + 0.57)
text(s, 0.9, 6.55, 11.5, 0.5,
     ("Every finding traces back to deterministic evidence — validated before it reaches the UI.",
      15, SLATE, False), align=PP_ALIGN.CENTER)

# 7 — Product experience (screenshot placeholders)
s = new()
eyebrow(s, 0.9, 0.6, "Product experience")
title(s, 0.9, 0.95, "One workspace, from signal to decision")
shots = ["Overview", "Investigation drawer", "Ask Pulse", "Executive brief"]
for i, sh in enumerate(shots):
    col, row = i % 2, i // 2
    x = 0.9 + col * 6.0
    y = 2.0 + row * 2.3
    box(s, x, y, 5.6, 2.0, fill=WHITE, line=BORDER)
    text(s, x + 0.2, y + 0.15, 5.2, 0.4, (sh, 15, NAVY, True))
    text(s, x + 0.2, y + 0.8, 5.2, 0.8,
         ("[ drop docs/screenshots/ image here ]", 12, SLATE, False), align=PP_ALIGN.CENTER)

# 8 — Business impact
s = new()
eyebrow(s, 0.9, 0.6, "Business impact")
title(s, 0.9, 0.95, "Faster, grounded, leadership-ready")
cards = [("Transport Manager", "Investigate in minutes, not hours — the evidence is pre-assembled.", BLUE),
         ("Facilities Head", "Leadership-ready view of cost, safety, reliability and risk.", TEAL),
         ("Organization", "Earlier detection of safety, billing and operational risks.", NAVY)]
for i, (h, body, col) in enumerate(cards):
    x = 0.9 + i * 3.95
    box(s, x, 2.3, 3.6, 2.8, fill=WHITE, line=BORDER)
    box(s, x, 2.3, 3.6, 0.12, fill=col, line=None, radius=False)
    text(s, x + 0.25, 2.6, 3.1, 0.6, (h, 17, col, True))
    text(s, x + 0.25, 3.3, 3.1, 1.6, (body, 14, SLATE, False))
text(s, 0.9, 5.6, 11.5, 0.5,
     ("Grounded in deterministic analytics — impact claims stay qualitative, no invented ROI.",
      13, SLATE, False), align=PP_ALIGN.CENTER)

# 9 — Enterprise evolution
s = new()
box(s, 0, 0, 13.333, 7.5, fill=NAVY, line=None, radius=False)
eyebrow(s, 0.9, 0.7, "Built for enterprise evolution", TEAL)
text(s, 0.9, 1.1, 11.5, 0.9, ("Today → Tomorrow", 32, WHITE, True))
box(s, 0.9, 2.3, 5.4, 3.2, fill=RGBColor(0x14,0x2F,0x52), line=None)
text(s, 1.2, 2.5, 4.9, 0.5, ("Today", 18, TEAL, True))
text(s, 1.2, 3.1, 4.9, 2.4,
     [("Hackathon prototype", 15, WHITE, True),
      ("Deterministic analytics + cross-domain engine", 13, RGBColor(0xB9,0xCD,0xF0), False),
      ("Embedded DuckDB, single tenant", 13, RGBColor(0xB9,0xCD,0xF0), False),
      ("Grounded, pluggable LLM narration", 13, RGBColor(0xB9,0xCD,0xF0), False)])
box(s, 6.9, 2.3, 5.4, 3.2, fill=RGBColor(0x14,0x2F,0x52), line=None)
text(s, 7.2, 2.5, 4.9, 0.5, ("Tomorrow", 18, TEAL, True))
text(s, 7.2, 3.1, 4.9, 2.4,
     [("Multi-tenant · real-time signals", 14, WHITE, False),
      ("Durable, governed action workflows", 14, WHITE, False),
      ("Audit trail · RBAC · model gateway", 14, WHITE, False),
      ("AWS data-lake deployment", 14, WHITE, False)])
text(s, 0.9, 5.9, 11.5, 0.9,
     ('"MoveInSync Pulse turns mobility data into decisions."', 22, TEAL, True),
     align=PP_ALIGN.CENTER)

out = Path(__file__).with_name("MoveInSync_Pulse_Hackathon.pptx")
prs.save(out)
print(f"Saved {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
